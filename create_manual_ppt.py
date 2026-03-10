"""
차병원 인턴 턴표 교환 시스템 - 교환 규칙 & 매뉴얼 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 (Teal Trust 기반 병원 테마) ──
C_DARK    = RGBColor(0x1A, 0x2E, 0x44)
C_TEAL    = RGBColor(0x02, 0x80, 0x90)
C_MINT    = RGBColor(0x02, 0xC3, 0x9A)
C_LIGHT   = RGBColor(0xF0, 0xF9, 0xF7)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x6B, 0x72, 0x80)
C_LGRAY   = RGBColor(0xE5, 0xE7, 0xEB)
C_RED     = RGBColor(0xEF, 0x44, 0x44)
C_ORANGE  = RGBColor(0xF9, 0x73, 0x16)
C_GREEN   = RGBColor(0x10, 0xB9, 0x81)
C_YELLOW  = RGBColor(0xFA, 0xCC, 0x15)
C_BG      = RGBColor(0xFA, 0xFB, 0xFC)
C_BLUE_BG = RGBColor(0xDB, 0xEA, 0xFE)
C_BLUE    = RGBColor(0x1D, 0x4E, 0xD8)
C_YELLOW_BG = RGBColor(0xFE, 0xF3, 0xC7)
C_GREEN_BG  = RGBColor(0xEC, 0xFD, 0xF5)
C_DARK_CARD = RGBColor(0x24, 0x3B, 0x55)

FONT = 'NanumGothic'

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

# ── 유틸리티 함수 ──
def add_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_rrect(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def T(slide, text, x, y, w, h, sz=14, color=C_DARK, bold=False, align=PP_ALIGN.LEFT):
    """단일 텍스트 추가"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    return txBox

def ML(slide, lines, x, y, w, h, sz=14, color=C_DARK, bold=False, align=PP_ALIGN.LEFT):
    """여러 줄 텍스트. lines: str | (text, bold, color)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sz * 0.4)
        p.space_before = Pt(0)
        run = p.add_run()
        if isinstance(line, tuple):
            run.text = line[0]
            run.font.bold = line[1] if len(line) > 1 else bold
            run.font.color.rgb = line[2] if len(line) > 2 else color
        else:
            run.text = line
            run.font.bold = bold
            run.font.color.rgb = color
        run.font.size = Pt(sz)
        run.font.name = FONT
    return txBox

def card(slide, icon, title, desc, x, y, w=4.0, h=1.2):
    add_rrect(slide, x, y, w, h, C_WHITE)
    T(slide, icon, x + 0.2, y + 0.15, 0.6, 0.6, sz=24, color=C_TEAL, align=PP_ALIGN.CENTER)
    T(slide, title, x + 0.85, y + 0.15, w - 1.1, 0.35, sz=13, color=C_DARK, bold=True)
    T(slide, desc, x + 0.85, y + 0.5, w - 1.1, h - 0.6, sz=10, color=C_GRAY)

def step(slide, num, x, y):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.4), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_TEAL
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(14)
    run.font.color.rgb = C_WHITE
    run.font.bold = True
    run.font.name = FONT

def header(slide, text):
    """공통 슬라이드 헤더 (진한 남색 바 + 흰 텍스트)"""
    add_bg(slide, C_BG)
    add_rect(slide, 0, 0, 10, 0.8, C_DARK)
    T(slide, text, 0.5, 0.15, 9, 0.5, sz=22, color=C_WHITE, bold=True)

def tip_bar(slide, text, y=4.7, bg=C_TEAL, color=C_WHITE):
    add_rrect(slide, 0.4, y, 9.2, 0.55, bg)
    T(slide, text, 0.6, y + 0.08, 8.8, 0.4, sz=12, color=color, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_DARK)
add_rect(s, 0, 0, 10, 0.08, C_TEAL)
add_rect(s, 0, 0.08, 10, 0.03, C_MINT)
T(s, "🏥", 4.25, 1.0, 1.5, 0.8, sz=40, color=C_WHITE, align=PP_ALIGN.CENTER)
T(s, "차병원 인턴 턴표 교환 시스템", 0.5, 1.9, 9, 0.7, sz=32, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
T(s, "교환 규칙 & 사용 매뉴얼", 0.5, 2.6, 9, 0.5, sz=18, color=C_MINT, align=PP_ALIGN.CENTER)
T(s, "2026년 인턴 교육과정", 0.5, 4.3, 9, 0.4, sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)
add_rect(s, 0, 5.545, 10, 0.08, C_TEAL)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_BG)
add_rect(s, 0, 0, 0.08, 5.625, C_TEAL)
T(s, "📋 목차", 0.5, 0.2, 9, 0.5, sz=28, color=C_DARK, bold=True)
add_rect(s, 0.5, 0.72, 2.5, 0.04, C_TEAL)

toc = [
    ("1", "로그인 & 화면 구성", "첫 접속, 비밀번호, 화면 설명"),
    ("2", "교환 규칙 한눈에 보기", "4가지 핵심 규칙 요약"),
    ("3", "필수과목 규칙", "IM, GS, OB, PE, 진로탐색"),
    ("4", "휴가턴 교환 규칙", "휴가턴 주고받기 밸런스"),
    ("5", "교환 신청 방법", "단일 교환 & 복합 교환"),
    ("6", "조합 찾기 & 장터", "자동 조합 탐색, 교환 장터"),
    ("7", "요청 수락/거절 & 시뮬레이션", "받은 요청 처리, 사전 탐색"),
]
for i, (num, title, desc) in enumerate(toc):
    yy = 1.0 + i * 0.6
    step(s, num, 0.6, yy)
    T(s, title, 1.15, yy, 4, 0.28, sz=14, color=C_DARK, bold=True)
    T(s, desc, 1.15, yy + 0.28, 4, 0.22, sz=9, color=C_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: 로그인 & 화면 구성 (NEW)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "1  로그인 & 화면 구성")

# 로그인 카드
add_rrect(s, 0.4, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 0.4, 1.0, 0.08, 2.0, C_TEAL)
T(s, "🔐 로그인", 0.7, 1.1, 4.0, 0.35, sz=15, color=C_TEAL, bold=True)
ML(s, [
    ("이름 + 비밀번호로 로그인", False, C_DARK),
    "",
    ("기본 비밀번호: 1234", True, C_DARK),
    "→ 첫 로그인 시 비밀번호 변경 화면이 나타나요",
    "→ 사이드바 > 🔑 비밀번호 변경에서 언제든 변경 가능",
], 0.7, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 화면 구성 카드
add_rrect(s, 5.2, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 5.2, 1.0, 0.08, 2.0, C_MINT)
T(s, "📱 화면 구성", 5.5, 1.1, 4.0, 0.35, sz=15, color=C_TEAL, bold=True)
ML(s, [
    ("메인 화면 (가운데)", True, C_DARK),
    "  교환 신청 / 전체 스케줄 / 교환 장터",
    "",
    ("사이드바 (왼쪽 > 버튼)", True, C_DARK),
    "  내 현황 / 받은 요청 / 교환 시뮬레이션",
], 5.5, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 하단: 전체 스케줄 보기
add_rrect(s, 0.4, 3.3, 9.2, 1.6, C_WHITE)
T(s, "📊 전체 스케줄 보기", 0.7, 3.4, 8.5, 0.35, sz=15, color=C_TEAL, bold=True)

# 왼쪽 설명
ML(s, [
    "모든 인턴의 턴표를 한눈에 볼 수 있어요",
    "",
    "노란 배경 = 내 행",
    "파란 배경 = 교환 신청한 상대방",
    "빨간 테두리 = 교환 중인 턴",
], 0.7, 3.8, 4.0, 1.0, sz=11, color=C_DARK)

# 오른쪽 팁
ML(s, [
    ("💡 화면 전환", True, C_TEAL),
    "오른쪽 상단에 📱/💻 버튼으로",
    "모바일 ↔ PC 레이아웃 전환 가능!",
    "",
    ("🔄 새로고침", True, C_TEAL),
    "다른 사람 교환 내용을 보려면",
    "🔄 새로고침 버튼을 눌러주세요",
], 5.5, 3.8, 3.9, 1.0, sz=10, color=C_DARK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: 교환 규칙 한눈에 보기
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "2  교환 규칙 한눈에 보기")

card(s, "📋", "필수과목 5개 유지",
     "교환 후에도 IM, GS, OB, PE, 진로탐색\n5개 과목이 반드시 모두 있어야 해요",
     0.4, 1.1, 4.3, 1.2)
card(s, "🚫", "1턴, 2턴 교환 불가",
     "1턴과 2턴은 교환할 수 없어요\n3턴~13턴만 교환 가능합니다",
     5.3, 1.1, 4.3, 1.2)
card(s, "🏖️", "휴가턴 교환 규칙",
     "휴가턴을 줬으면 같은 상대에게서\n다른 휴가턴을 받아와야 해요",
     0.4, 2.6, 4.3, 1.2)
card(s, "🏥", "분당 근무 최소 7턴",
     "교환 후에도 분당 근무가\n최소 7턴 이상이어야 해요",
     5.3, 2.6, 4.3, 1.2)

tip_bar(s, "💡 모든 규칙은 시스템이 자동으로 검증해요! 규칙 위반 교환은 신청 자체가 불가능합니다", y=4.1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: 필수과목 규칙
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "3  필수과목 규칙")

depts = [("IM", "내과"), ("GS", "외과"), ("OB", "산부인과"), ("PE", "소아과"), ("진로탐색", "진로탐색")]
for i, (code, name) in enumerate(depts):
    xx = 0.4 + i * 1.9
    add_rrect(s, xx, 1.1, 1.7, 1.0, C_WHITE)
    T(s, code, xx, 1.2, 1.7, 0.45, sz=18, color=C_TEAL, bold=True, align=PP_ALIGN.CENTER)
    T(s, name, xx, 1.65, 1.7, 0.3, sz=10, color=C_GRAY, align=PP_ALIGN.CENTER)

ML(s, [
    ("교환 후에도 위 5개 필수과목이 모두 스케줄에 있어야 합니다.", True, C_DARK),
    "",
    "예) 내가 IM을 주고 GS를 받으면, 내 스케줄에 IM이 없어지죠?",
    "     → 다른 턴에 IM이 이미 있으면 OK!",
    "     → IM이 아예 사라지면 교환 불가 (시스템이 자동 차단)",
], 0.6, 2.4, 8.8, 2.5, sz=13, color=C_DARK)

tip_bar(s, "💡 Tip: 조건 불충족 시 '가능한 조합 찾기'를 누르면 추가 교환으로 해결되는 조합을 자동 탐색해줘요!",
        y=4.4, bg=C_YELLOW_BG, color=C_ORANGE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: 휴가턴 교환 규칙 (핵심!)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "4  휴가턴 교환 규칙  ⭐핵심")

# 단일 교환 (불가)
add_rrect(s, 0.4, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 0.4, 1.0, 0.08, 2.0, C_RED)
T(s, "❌ 단일 교환 — 한쪽만 휴가면 불가", 0.7, 1.1, 4.0, 0.35, sz=13, color=C_RED, bold=True)
ML(s, [
    "A의 휴가턴 ↔ B의 일반턴",
    "",
    "→ A가 휴가를 잃고 돌려받을 방법이 없음",
    "→ 시스템이 자동 차단합니다",
], 0.7, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 복합 교환 (OK)
add_rrect(s, 5.2, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 5.2, 1.0, 0.08, 2.0, C_GREEN)
T(s, "✅ 복합 교환 — 수지만 맞으면 OK", 5.5, 1.1, 4.0, 0.35, sz=13, color=C_GREEN, bold=True)
ML(s, [
    "A가 B에게 휴가턴 1개 줌",
    "+ B에게서 다른 휴가턴 1개 받음",
    "",
    "→ 상대방별 휴가 수지가 같으면 통과!",
], 5.5, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 예시
add_rrect(s, 0.4, 3.2, 9.2, 1.8, C_GREEN_BG)
T(s, "📌 예시로 이해하기", 0.7, 3.3, 8.5, 0.35, sz=14, color=C_TEAL, bold=True)
ML(s, [
    ("김인턴의 5턴(휴가)을 이인턴에게 주고 싶어요", False, C_DARK),
    "",
    ("→ 방법: 5턴 교환 + 9턴도 함께 교환 (이인턴의 9턴이 휴가)", False, C_DARK),
    ("   김인턴이 5턴 휴가를 주고(→준 1개), 9턴에서 이인턴 휴가를 받음(→받은 1개)", False, C_TEAL),
    ("   수지: 준 1개 = 받은 1개 ✅  복합 교환으로 한 번에 신청!", True, C_GREEN),
], 0.7, 3.65, 8.5, 1.3, sz=11)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: 교환 신청 방법
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "5  교환 신청 방법")

# Step 1
step(s, 1, 0.5, 1.1)
T(s, "상대방 & 턴 선택", 1.1, 1.1, 3.5, 0.35, sz=15, color=C_DARK, bold=True)
T(s, "교환하고 싶은 상대와 턴을 선택하면\n\"나: ○○ ↔ 상대: ○○\" 미리보기가 표시돼요",
   1.1, 1.5, 3.8, 0.6, sz=11, color=C_GRAY)

# Step 2
step(s, 2, 0.5, 2.3)
T(s, "조건 확인", 1.1, 2.3, 3.5, 0.35, sz=15, color=C_DARK, bold=True)
ML(s, [
    ("✅ 초록색 = 교환 가능!", True, C_GREEN),
    ("⚠️ 빨간색 = 규칙 위반 (항목 추가 or 조합 찾기)", True, C_RED),
], 1.1, 2.7, 3.8, 0.6, sz=11)

# Step 3
step(s, 3, 0.5, 3.5)
T(s, "요청 보내기", 1.1, 3.5, 3.5, 0.35, sz=15, color=C_DARK, bold=True)
T(s, "📨 요청 보내기 버튼을 누르면\n상대방에게 교환 요청이 전달돼요",
   1.1, 3.9, 3.8, 0.6, sz=11, color=C_GRAY)

# 오른쪽: 복합 교환
add_rrect(s, 5.2, 1.0, 4.4, 2.5, C_WHITE)
T(s, "🔗 복합 교환 (여러 건 동시)", 5.5, 1.1, 3.9, 0.35, sz=14, color=C_TEAL, bold=True)
ML(s, [
    "➕ 항목 추가 버튼으로 교환을 추가",
    "",
    "• 2~3개 턴을 한 번에 교환 신청",
    "• 모든 상대방이 수락해야 실행",
    "• 하나라도 거절하면 전체 취소",
    "",
    "💡 휴가턴 교환이나 과목 규칙 해결에",
    "   복합 교환이 꼭 필요해요!",
], 5.5, 1.5, 3.9, 2.0, sz=11, color=C_DARK)

# 하단 팁
add_rrect(s, 5.2, 3.7, 4.4, 1.0, C_BLUE_BG)
T(s, "🔍 조건 불충족 시", 5.5, 3.8, 3.9, 0.3, sz=12, color=C_BLUE, bold=True)
ML(s, [
    "항목 추가: 수동으로 교환 추가",
    "가능한 조합 찾기: 자동으로 조합 탐색!",
], 5.5, 4.1, 3.9, 0.5, sz=10, color=C_BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: 가능한 조합 찾기 & 교환 장터
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "6  조합 찾기 & 교환 장터")

# 왼쪽: 조합 찾기
add_rrect(s, 0.4, 1.0, 4.4, 3.5, C_WHITE)
T(s, "🔍 가능한 조합 찾기", 0.7, 1.1, 4.0, 0.4, sz=16, color=C_TEAL, bold=True)
ML(s, [
    "교환이 규칙에 안 맞을 때 자동으로",
    "추가 교환 조합을 탐색해줘요!",
    "",
    ("사용법:", True, C_DARK),
    "1. 원하는 교환을 선택",
    "2. 조건 불충족 시 🔍 버튼 클릭",
    "3. 가능한 조합 목록에서 선택",
    "4. ⬆️ 설정 버튼으로 자동 세팅",
    "",
    ("💡 시스템이 최대 200개 조합을", False, C_TEAL),
    ("   자동으로 탐색해줍니다!", False, C_TEAL),
], 0.7, 1.55, 3.9, 2.8, sz=11, color=C_DARK)

# 오른쪽: 교환 장터
add_rrect(s, 5.2, 1.0, 4.4, 3.5, C_WHITE)
T(s, "🏪 교환 장터", 5.5, 1.1, 4.0, 0.4, sz=16, color=C_ORANGE, bold=True)
ML(s, [
    "직접 상대를 찾기 어려울 때",
    "게시판에 교환을 올려보세요!",
    "",
    ("장터 둘러보기:", True, C_DARK),
    "🟢 초록 점 = 바로 교환 가능",
    "🔴 빨간 점 = 복합 교환으로 가능할 수도",
    "",
    ("내 턴 올리기:", True, C_DARK),
    "• 줄 턴 지정: \"5턴을 줄게, IM 받고 싶어\"",
    "• 받을 턴 지정: \"GS 받고 싶어\"",
    "• 메시지도 함께 남길 수 있어요",
], 5.5, 1.55, 3.9, 2.8, sz=11, color=C_DARK)

tip_bar(s, "💡 장터에서 '응하기' 버튼을 누르면 바로 교환 신청으로 연결돼요!")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: 요청 수락/거절 & 시뮬레이션
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, "7  요청 수락 / 거절 & 시뮬레이션")

# 수락
add_rrect(s, 0.4, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 0.4, 1.0, 0.08, 2.0, C_GREEN)
T(s, "✅ 수락하기", 0.7, 1.1, 4.0, 0.35, sz=15, color=C_GREEN, bold=True)
ML(s, [
    "• 사이드바에서 받은 요청 확인",
    "• ✅ 수락 버튼 → 즉시 교환 실행",
    "• 구글 시트에 자동 반영!",
    "",
    "⚠️ 복합 교환은 모든 사람이",
    "   수락해야 최종 실행돼요",
], 0.7, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 거절
add_rrect(s, 5.2, 1.0, 4.4, 2.0, C_WHITE)
add_rect(s, 5.2, 1.0, 0.08, 2.0, C_RED)
T(s, "❌ 거절하기", 5.5, 1.1, 4.0, 0.35, sz=15, color=C_RED, bold=True)
ML(s, [
    "• ❌ 거절 버튼으로 거절",
    "• 복합 교환은 하나만 거절해도",
    "  전체가 자동 취소됩니다",
    "",
    "💡 새로운 요청이 오면",
    "   사이드바에 숫자로 알려줘요",
], 5.5, 1.5, 3.9, 1.4, sz=11, color=C_DARK)

# 시뮬레이션
add_rrect(s, 0.4, 3.3, 9.2, 1.7, C_WHITE)
T(s, "🧪 교환 시뮬레이션 (사이드바 하단)", 0.7, 3.4, 8.5, 0.35, sz=14, color=C_TEAL, bold=True)
ML(s, [
    ("🔄 특정 턴 교환", True, C_DARK),
    "   특정 턴에서 교환 가능한 파트너 전체 목록",
    ("🎯 특정 턴 받기", True, C_DARK),
    "   내가 원하는 과목을 받을 수 있는 모든 조합 탐색",
    ("🔗 복합 교환 탐색", True, C_DARK),
    "   2~3개 턴을 동시에 교환해야만 가능한 조합 검색",
    "",
    "→ 요청 버튼을 누르면 바로 교환 신청 화면으로 연결돼요!",
], 0.7, 3.8, 8.5, 1.2, sz=10, color=C_DARK)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: 마무리
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, C_DARK)
add_rect(s, 0, 0, 10, 0.08, C_TEAL)
add_rect(s, 0, 0.08, 10, 0.03, C_MINT)

T(s, "🏥", 4.25, 1.1, 1.5, 0.8, sz=40, color=C_WHITE, align=PP_ALIGN.CENTER)
T(s, "준비 완료!", 0.5, 2.0, 9, 0.6, sz=30, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
T(s, "이제 턴표 교환을 시작해보세요", 0.5, 2.6, 9, 0.5, sz=16, color=C_MINT, align=PP_ALIGN.CENTER)

info = [
    ("🔐", "기본 비밀번호: 1234\n(첫 로그인 시 변경 필수!)"),
    ("📱", "모바일에서도 접속 가능\n(오른쪽 상단에서 화면 전환)"),
    ("🔄", "다른 사람 교환 확인은\n새로고침 버튼 클릭"),
]
for i, (icon, text) in enumerate(info):
    xx = 0.7 + i * 3.1
    add_rrect(s, xx, 3.3, 2.8, 1.2, C_DARK_CARD)
    T(s, icon, xx + 0.1, 3.4, 0.5, 0.5, sz=20, color=C_WHITE, align=PP_ALIGN.CENTER)
    T(s, text, xx + 0.6, 3.4, 2.0, 1.0, sz=10, color=C_LGRAY)

add_rect(s, 0, 5.545, 10, 0.08, C_TEAL)


# ── 저장 ──
output_path = r"C:\change\차병원_턴표교환_매뉴얼.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
